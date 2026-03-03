"""
Content Processing Utilities
Extracts text from YouTube, PDF, PPT, and Word documents
Supports multilingual YouTube videos (Hindi, Marathi, etc.) with auto-translation
"""
import os
import re
from typing import List, Dict, Tuple


# Supported languages for YouTube transcript extraction (priority order)
SUPPORTED_LANGUAGES = ['en', 'hi', 'mr', 'ta', 'te', 'kn', 'ml', 'gu', 'bn', 'pa', 'ur']
LANGUAGE_NAMES = {
    'en': 'English', 'hi': 'Hindi', 'mr': 'Marathi', 'ta': 'Tamil',
    'te': 'Telugu', 'kn': 'Kannada', 'ml': 'Malayalam', 'gu': 'Gujarati',
    'bn': 'Bengali', 'pa': 'Punjabi', 'ur': 'Urdu'
}


def process_content(content):
    """
    Main content processing function - SKIP transcript extraction for YouTube
    We use workspace name for question generation instead
    """
    try:
        if content.content_type == 'youtube':
            # SKIP transcript extraction - we use workspace name instead
            print(f"[Content] Skipping transcript extraction for YouTube video (using workspace name)")
            content.transcript = ""
            content.source_language = 'en'
        elif content.content_type == 'pdf':
            content.transcript = extract_pdf_text(content.file.path)
            content.source_language = 'en'
        elif content.content_type == 'ppt':
            content.transcript = extract_ppt_text(content.file.path)
            content.source_language = 'en'
        elif content.content_type == 'word':
            content.transcript = extract_word_text(content.file.path)
            content.source_language = 'en'
        
        # Skip key concepts extraction for YouTube
        if content.content_type != 'youtube' and content.transcript and len(content.transcript) > 50:
            content.key_concepts = extract_key_concepts(content.transcript)
        else:
            content.key_concepts = []
        
        content.processed = True
        content.save()
        
        return content
    except Exception as e:
        print(f"[Content] Processing error: {e}")
        # Mark as processed anyway so it doesn't block
        content.processed = True
        content.transcript = ""
        content.save()
        return content


def _retry_with_backoff(func, max_attempts=2, base_delay=2.0, label=""):
    """
    Retry a function with exponential backoff.
    Returns (success, result_or_error).
    """
    import time
    import random
    
    last_error = None
    for attempt in range(1, max_attempts + 1):
        try:
            result = func()
            return True, result
        except Exception as e:
            last_error = e
            if attempt < max_attempts:
                delay = base_delay * attempt + random.uniform(0, 1)
                print(f"[Transcript] {label} attempt {attempt} failed: {str(e)[:120]}. Retrying in {delay:.1f}s...")
                time.sleep(delay)
            else:
                print(f"[Transcript] {label} attempt {attempt} failed (final): {str(e)[:200]}")
    return False, last_error


def extract_youtube_transcript(url: str) -> Tuple[str, str]:
    """
    Extract transcript from YouTube video with 4-layer bulletproof fallback.
    
    Fallback chain:
        1. youtube-transcript-api (fast, reliable for most videos)
        2. yt-dlp subtitle extraction (different approach, handles some edge cases)
        3. Raw Innertube API scraping (no library dependency, bypasses library bugs)
        4. yt-dlp audio download + Gemini AI transcription (works even with NO captions)
    
    Each method is retried up to 2 times with backoff before falling through.
    
    Returns:
        Tuple of (transcript_text, source_language_code)
    """
    import time
    import random
    
    # Check if it's a playlist
    if 'playlist' in url or 'list=' in url:
        return extract_youtube_playlist_transcript(url)
    
    # Extract video ID from URL
    video_id = extract_video_id(url)
    
    if not video_id:
        raise ValueError("Invalid YouTube URL")
    
    errors = []
    
    # ── Method 1: youtube-transcript-api ──
    print(f"[Transcript] Method 1: youtube-transcript-api for {video_id}")
    time.sleep(random.uniform(0.3, 1.0))
    ok, result = _retry_with_backoff(
        lambda: _extract_transcript_multilingual(video_id),
        max_attempts=2, label="youtube-transcript-api"
    )
    if ok:
        text, lang = result
        if text and len(text.strip()) > 30:
            print(f"[Transcript] ✅ Method 1 succeeded ({len(text)} chars, lang={lang})")
            return text, lang
        else:
            errors.append("youtube-transcript-api returned empty/short transcript")
    else:
        errors.append(f"youtube-transcript-api: {str(result)[:200]}")
    
    # ── Method 2: yt-dlp subtitles ──
    print(f"[Transcript] Method 2: yt-dlp subtitles for {video_id}")
    ok, result = _retry_with_backoff(
        lambda: _extract_transcript_ytdlp(video_id),
        max_attempts=2, label="yt-dlp-subtitles"
    )
    if ok:
        text, lang = result
        if text and len(text.strip()) > 30:
            print(f"[Transcript] ✅ Method 2 succeeded ({len(text)} chars, lang={lang})")
            return text, lang
        else:
            errors.append("yt-dlp returned empty/short subtitle text")
    else:
        errors.append(f"yt-dlp: {str(result)[:200]}")
    
    # ── Method 3: Raw Innertube API scraping ──
    print(f"[Transcript] Method 3: Raw Innertube API for {video_id}")
    ok, result = _retry_with_backoff(
        lambda: _extract_transcript_innertube(video_id),
        max_attempts=2, label="innertube-raw"
    )
    if ok:
        text, lang = result
        if text and len(text.strip()) > 30:
            print(f"[Transcript] ✅ Method 3 succeeded ({len(text)} chars, lang={lang})")
            return text, lang
        else:
            errors.append("Innertube returned empty/short transcript")
    else:
        errors.append(f"innertube: {str(result)[:200]}")
    
    # ── Method 4: yt-dlp audio + Gemini AI transcription ──
    print(f"[Transcript] Method 4: Audio + Gemini AI transcription for {video_id}")
    ok, result = _retry_with_backoff(
        lambda: _extract_transcript_audio_gemini(video_id),
        max_attempts=2, label="audio-gemini"
    )
    if ok:
        text, lang = result
        if text and len(text.strip()) > 30:
            print(f"[Transcript] ✅ Method 4 (Gemini AI) succeeded ({len(text)} chars)")
            return text, lang
        else:
            errors.append("Gemini AI returned empty/short transcript")
    else:
        errors.append(f"audio-gemini: {str(result)[:200]}")
    
    # All 4 methods failed
    error_summary = " | ".join(errors)
    raise ValueError(
        f"All 4 transcript methods failed for video {video_id}. "
        f"Errors: {error_summary}"
    )


def _extract_transcript_ytdlp(video_id: str) -> Tuple[str, str]:
    """
    Extract transcript using yt-dlp with proxy support.
    
    Returns:
        Tuple of (transcript_text, language_code)
    """
    import yt_dlp
    
    url = f"https://www.youtube.com/watch?v={video_id}"
    
    # Check for proxy
    proxy = os.environ.get('HTTP_PROXY') or os.environ.get('HTTPS_PROXY')
    
    # Look for cookies file
    cookies_path = os.path.join(os.path.dirname(__file__), 'youtube_cookies.txt')
    
    # Configure yt-dlp with minimal options to avoid format errors
    ydl_opts = {
        'skip_download': True,
        'writesubtitles': True,
        'writeautomaticsub': True,
        'subtitleslangs': ['en', 'en-US', 'en-GB'],
        'quiet': False,  # Show errors for debugging
        'no_warnings': False,
        'extract_flat': False,
        'ignoreerrors': False,
        # Do NOT specify 'format' - we only want subtitles, not media
    }
    
    # Add cookies if file exists
    if os.path.exists(cookies_path):
        ydl_opts['cookiefile'] = cookies_path
        print(f"[Transcript] yt-dlp using cookies from {cookies_path}")
    
    if proxy:
        print(f"[Transcript] yt-dlp using proxy: {proxy}")
        ydl_opts['proxy'] = proxy
    
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            
            # Get subtitles
            subtitles = info.get('subtitles', {})
            automatic_captions = info.get('automatic_captions', {})
            
            print(f"[Transcript] Available subtitle languages: {list(subtitles.keys())}")
            print(f"[Transcript] Available auto-caption languages: {list(automatic_captions.keys())}")
            
            # Try English variants
            sub_data = None
            for lang in ['en', 'en-US', 'en-GB', 'en-CA']:
                sub_data = subtitles.get(lang) or automatic_captions.get(lang)
                if sub_data:
                    print(f"[Transcript] Using language: {lang}")
                    break
            
            if not sub_data:
                raise ValueError("No English subtitles found")
            
            # Get subtitle URL (prefer vtt or srv3 format)
            sub_url = None
            for fmt in sub_data:
                if fmt.get('ext') in ['vtt', 'srv3', 'srv2', 'srv1', 'json3']:
                    sub_url = fmt['url']
                    break
            
            if not sub_url and sub_data:
                sub_url = sub_data[0]['url']
            
            if not sub_url:
                raise ValueError("No subtitle URL found")
            
            print(f"[Transcript] Downloading subtitles from: {sub_url[:100]}...")
            
            # Download and parse subtitles with proxy
            import requests
            proxies = {'http': proxy, 'https': proxy} if proxy else None
            response = requests.get(sub_url, timeout=10, proxies=proxies)
            content = response.text
            
            # Parse VTT format
            if 'WEBVTT' in content:
                lines = content.split('\n')
                text_parts = []
                for line in lines:
                    line = line.strip()
                    # Skip timestamps and metadata
                    if line and not line.startswith('WEBVTT') and not '-->' in line and not line.isdigit():
                        # Remove HTML tags
                        import re
                        clean_line = re.sub(r'<[^>]+>', '', line)
                        if clean_line:
                            text_parts.append(clean_line)
                
                text = ' '.join(text_parts)
                print(f"[Transcript] yt-dlp extracted transcript ({len(text)} chars)")
                return text, 'en'
            
            # Parse JSON format (json3)
            elif content.startswith('{') or content.strip().startswith('{'):
                import json
                data = json.loads(content)
                text_parts = []
                
                # Handle different JSON structures
                if 'events' in data:
                    for event in data['events']:
                        if 'segs' in event:
                            for seg in event['segs']:
                                if 'utf8' in seg:
                                    text_parts.append(seg['utf8'])
                
                text = ' '.join(text_parts)
                if text:
                    print(f"[Transcript] yt-dlp extracted transcript ({len(text)} chars)")
                    return text, 'en'
            
            raise ValueError("Could not parse subtitle format")
    
    except Exception as e:
        error_msg = str(e)
        print(f"[Transcript] yt-dlp error details: {error_msg}")
        raise ValueError(f"yt-dlp extraction failed: {error_msg}")


def _extract_transcript_innertube(video_id: str) -> Tuple[str, str]:
    """
    Method 3: Extract transcript by scraping YouTube's Innertube API directly.
    
    This bypasses both youtube-transcript-api and yt-dlp by:
    1. Fetching the YouTube watch page HTML
    2. Extracting ytInitialPlayerResponse JSON
    3. Finding caption track URLs
    4. Downloading and parsing the caption XML
    
    Returns:
        Tuple of (transcript_text, language_code)
    """
    import requests
    import json
    import xml.etree.ElementTree as ET
    
    url = f"https://www.youtube.com/watch?v={video_id}"
    
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
        'Accept-Language': 'en-US,en;q=0.9',
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
    }
    
    proxy = os.environ.get('HTTP_PROXY') or os.environ.get('HTTPS_PROXY')
    proxies = {'http': proxy, 'https': proxy} if proxy else None
    
    # Step 1: Fetch the watch page
    response = requests.get(url, headers=headers, timeout=15, proxies=proxies)
    response.raise_for_status()
    html = response.text
    
    # Step 2: Extract ytInitialPlayerResponse
    player_response = None
    
    # Try pattern 1: var ytInitialPlayerResponse = {...};
    match = re.search(r'var\s+ytInitialPlayerResponse\s*=\s*(\{.+?\})\s*;', html)
    if match:
        try:
            player_response = json.loads(match.group(1))
        except json.JSONDecodeError:
            pass
    
    # Try pattern 2: ytInitialPlayerResponse inside ytInitialData or script
    if not player_response:
        match = re.search(r'ytInitialPlayerResponse\s*=\s*(\{.+?\})\s*;', html)
        if match:
            try:
                player_response = json.loads(match.group(1))
            except json.JSONDecodeError:
                pass
    
    # Try pattern 3: look in inline script with playerResponse
    if not player_response:
        match = re.search(r'"playerResponse"\s*:\s*"(.+?)"', html)
        if match:
            try:
                decoded = match.group(1).encode().decode('unicode_escape')
                player_response = json.loads(decoded)
            except (json.JSONDecodeError, UnicodeDecodeError):
                pass
    
    if not player_response:
        raise ValueError("Could not extract ytInitialPlayerResponse from page")
    
    # Step 3: Find caption tracks
    captions = player_response.get('captions', {})
    tracklist = captions.get('playerCaptionsTracklistRenderer', {})
    caption_tracks = tracklist.get('captionTracks', [])
    
    if not caption_tracks:
        raise ValueError("No caption tracks found in player response")
    
    print(f"[Transcript] Innertube found {len(caption_tracks)} caption tracks")
    
    # Prefer English, then any supported language
    selected_track = None
    for lang_code in ['en', 'en-US', 'en-GB'] + SUPPORTED_LANGUAGES:
        for track in caption_tracks:
            if track.get('languageCode', '').startswith(lang_code.split('-')[0]):
                selected_track = track
                break
        if selected_track:
            break
    
    if not selected_track:
        # Use any available track
        selected_track = caption_tracks[0]
    
    caption_url = selected_track.get('baseUrl', '')
    if not caption_url:
        raise ValueError("Caption track has no baseUrl")
    
    track_lang = selected_track.get('languageCode', 'en')
    print(f"[Transcript] Innertube downloading captions (lang={track_lang})")
    
    # Step 4: Download caption XML
    caption_response = requests.get(caption_url, headers=headers, timeout=10, proxies=proxies)
    caption_response.raise_for_status()
    
    # Parse XML captions
    try:
        root = ET.fromstring(caption_response.text)
        text_parts = []
        for elem in root.iter('text'):
            text_content = elem.text
            if text_content:
                # Unescape HTML entities
                text_content = text_content.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&#39;', "'").replace('&quot;', '"')
                text_parts.append(text_content.strip())
        
        text = ' '.join(text_parts)
        if text:
            print(f"[Transcript] Innertube extracted {len(text)} chars")
            return text, track_lang
        raise ValueError("Caption XML had no text content")
    except ET.ParseError:
        # Maybe it's not XML, try as plain text
        text = re.sub(r'<[^>]+>', ' ', caption_response.text)
        text = ' '.join(text.split())
        if len(text) > 30:
            return text, track_lang
        raise ValueError("Could not parse caption data")


def _extract_transcript_audio_gemini(video_id: str) -> Tuple[str, str]:
    """
    Method 4 (Ultimate Fallback): Download audio and use Gemini AI to transcribe.
    
    This works even when a video has NO captions/subtitles at all.
    Uses yt-dlp to download the audio stream, then sends it to Gemini
    for transcription.
    
    Returns:
        Tuple of (transcript_text, language_code)
    """
    import tempfile
    import yt_dlp
    
    try:
        from django.conf import settings
        api_key = getattr(settings, 'GEMINI_API_KEY', None)
    except Exception:
        api_key = None
    
    if not api_key:
        api_key = os.environ.get('GEMINI_API_KEY')
    
    if not api_key:
        raise ValueError("No GEMINI_API_KEY configured — cannot use AI transcription fallback")
    
    url = f"https://www.youtube.com/watch?v={video_id}"
    
    # Download audio to temp file
    with tempfile.TemporaryDirectory() as tmpdir:
        audio_path = os.path.join(tmpdir, 'audio.m4a')
        
        proxy = os.environ.get('HTTP_PROXY') or os.environ.get('HTTPS_PROXY')
        cookies_path = os.path.join(os.path.dirname(__file__), 'youtube_cookies.txt')
        
        ydl_opts = {
            'format': 'bestaudio[ext=m4a]/bestaudio/best',
            'outtmpl': audio_path,
            'quiet': True,
            'no_warnings': True,
            'extract_flat': False,
            # Limit to 15 mins of audio to stay within Gemini limits
            'download_ranges': None,
            'postprocessors': [],
        }
        
        if os.path.exists(cookies_path):
            ydl_opts['cookiefile'] = cookies_path
        if proxy:
            ydl_opts['proxy'] = proxy
        
        print(f"[Transcript] Downloading audio for {video_id}...")
        
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        
        # Find the downloaded file (yt-dlp may add extension)
        actual_audio = None
        for f in os.listdir(tmpdir):
            if f.startswith('audio'):
                actual_audio = os.path.join(tmpdir, f)
                break
        
        if not actual_audio or not os.path.exists(actual_audio):
            raise ValueError("Audio download failed - no file produced")
        
        file_size = os.path.getsize(actual_audio)
        print(f"[Transcript] Audio downloaded ({file_size / 1024 / 1024:.1f} MB)")
        
        # Check file size (Gemini inline limit is ~20MB)
        if file_size > 20 * 1024 * 1024:
            print(f"[Transcript] Audio too large for inline upload, using File API...")
            return _transcribe_with_gemini_file_api(actual_audio, api_key)
        
        # Use Gemini to transcribe
        return _transcribe_with_gemini_inline(actual_audio, api_key)


def _transcribe_with_gemini_inline(audio_path: str, api_key: str) -> Tuple[str, str]:
    """Transcribe audio using Gemini inline upload."""
    import google.generativeai as genai
    
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-2.0-flash')
    
    # Read audio file
    with open(audio_path, 'rb') as f:
        audio_data = f.read()
    
    # Determine MIME type
    ext = os.path.splitext(audio_path)[1].lower()
    mime_map = {
        '.m4a': 'audio/mp4',
        '.mp3': 'audio/mpeg',
        '.wav': 'audio/wav',
        '.ogg': 'audio/ogg',
        '.webm': 'audio/webm',
    }
    mime_type = mime_map.get(ext, 'audio/mp4')
    
    print(f"[Transcript] Sending audio to Gemini ({len(audio_data) / 1024 / 1024:.1f} MB, {mime_type})...")
    
    response = model.generate_content([
        {
            'mime_type': mime_type,
            'data': audio_data
        },
        "Transcribe this audio completely and accurately. "
        "Output ONLY the transcription text, nothing else. "
        "If the audio is not in English, translate it to English. "
        "Do not add timestamps or speaker labels."
    ])
    
    text = response.text.strip()
    if text:
        print(f"[Transcript] Gemini transcription: {len(text)} chars")
        return text, 'en'
    
    raise ValueError("Gemini returned empty transcription")


def _transcribe_with_gemini_file_api(audio_path: str, api_key: str) -> Tuple[str, str]:
    """Transcribe large audio files using Gemini File API."""
    import google.generativeai as genai
    import time
    
    genai.configure(api_key=api_key)
    
    print(f"[Transcript] Uploading audio via File API...")
    audio_file = genai.upload_file(audio_path)
    
    # Wait for processing
    while audio_file.state.name == "PROCESSING":
        print(f"[Transcript] File processing...")
        time.sleep(3)
        audio_file = genai.get_file(audio_file.name)
    
    if audio_file.state.name == "FAILED":
        raise ValueError(f"Gemini file processing failed: {audio_file.state}")
    
    model = genai.GenerativeModel('gemini-2.0-flash')
    
    response = model.generate_content([
        audio_file,
        "Transcribe this audio completely and accurately. "
        "Output ONLY the transcription text, nothing else. "
        "If the audio is not in English, translate it to English. "
        "Do not add timestamps or speaker labels."
    ])
    
    # Clean up uploaded file
    try:
        genai.delete_file(audio_file.name)
    except Exception:
        pass
    
    text = response.text.strip()
    if text:
        print(f"[Transcript] Gemini File API transcription: {len(text)} chars")
        return text, 'en'
    
    raise ValueError("Gemini File API returned empty transcription")


def _extract_transcript_multilingual(video_id: str) -> Tuple[str, str]:
    """
    Extract transcript from YouTube video with cookie support.
    Compatible with youtube-transcript-api v1.x+
    
    Returns:
        Tuple of (english_transcript, original_language_code)
    """
    from youtube_transcript_api import YouTubeTranscriptApi
    from http.cookiejar import MozillaCookieJar
    import requests
    
    def _snippets_to_text(snippets):
        """Convert transcript snippets to plain text."""
        parts = []
        for s in snippets:
            if isinstance(s, dict):
                parts.append(s.get('text', ''))
            elif hasattr(s, 'text'):
                # v1.2+ uses FetchedTranscriptSnippet objects with .text attribute
                parts.append(s.text)
            else:
                parts.append(str(s))
        return ' '.join(parts)
    
    # Look for cookies file
    cookies_path = os.path.join(os.path.dirname(__file__), 'youtube_cookies.txt')
    
    # Create session with cookies and headers
    session = requests.Session()
    
    # Load cookies if available
    if os.path.exists(cookies_path):
        try:
            cookie_jar = MozillaCookieJar(cookies_path)
            cookie_jar.load(ignore_discard=True, ignore_expires=True)
            session.cookies = cookie_jar
            print(f"[Transcript] Using cookies from {cookies_path}")
        except Exception as e:
            print(f"[Transcript] Could not load cookies: {e}")
    
    # Add headers
    session.headers.update({
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
        'Accept-Language': 'en-US,en;q=0.9',
    })
    
    # Create API instance — v1.2+ accepts http_client in constructor for auth/cookies
    api = YouTubeTranscriptApi(http_client=session)
    
    # Try languages in priority order
    for lang_code in SUPPORTED_LANGUAGES:
        try:
            transcript_list = api.list(video_id)
            transcript = transcript_list.find_transcript([lang_code])
            data = transcript.fetch()
            text = _snippets_to_text(data)
            lang_name = LANGUAGE_NAMES.get(lang_code, lang_code)
            print(f"[Transcript] Found {lang_name} transcript for {video_id}")
            return text, lang_code
        except Exception:
            continue
    
    # If no supported language found, try any available
    try:
        transcript_list = api.list(video_id)
        transcript = list(transcript_list)[0]
        data = transcript.fetch()
        text = _snippets_to_text(data)
        print(f"[Transcript] Found transcript for {video_id}")
        return text, transcript.language_code
    except Exception as e:
        raise ValueError(f"No transcripts/captions found: {str(e)}")


def _translate_via_grok(text: str, source_lang: str) -> str:
    """
    Translate text to English using Grok AI API.
    Used as fallback when YouTube auto-translate is not available.
    
    Args:
        text: Text in source language
        source_lang: Source language code (hi, mr, etc.)
    
    Returns:
        Translated English text, or None if translation fails
    """
    try:
        import requests
        
        api_key = os.environ.get('XAI_API_KEY')
        if not api_key:
            print("[Translation] No XAI_API_KEY found, cannot translate via Grok AI")
            return None
        
        lang_name = LANGUAGE_NAMES.get(source_lang, source_lang)
        
        # Limit text to avoid token limits (keep first ~4000 chars for translation, do in chunks if needed)
        max_chunk_size = 4000
        
        if len(text) <= max_chunk_size:
            chunks = [text]
        else:
            # Split into chunks at sentence boundaries
            chunks = []
            current_chunk = ""
            sentences = text.replace('। ', '।\n').replace('. ', '.\n').split('\n')
            for sentence in sentences:
                if len(current_chunk) + len(sentence) < max_chunk_size:
                    current_chunk += sentence + " "
                else:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = sentence + " "
            if current_chunk:
                chunks.append(current_chunk.strip())
        
        translated_chunks = []
        
        for i, chunk in enumerate(chunks):
            prompt = f"""Translate the following {lang_name} text to English accurately. 
Preserve the original meaning, technical terms, and context. 
Return ONLY the English translation, nothing else.

{lang_name} text:
{chunk}"""
            
            response = requests.post(
                "https://api.x.ai/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "grok-beta",
                    "messages": [
                        {"role": "system", "content": f"You are an expert {lang_name}-to-English translator. Translate accurately, preserving technical terms and educational context. Return only the translation."},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.3,
                    "max_tokens": 3000
                },
                timeout=30
            )
            
            response.raise_for_status()
            result = response.json()
            translated_text = result['choices'][0]['message']['content'].strip()
            translated_chunks.append(translated_text)
            print(f"[Translation] Chunk {i+1}/{len(chunks)} translated successfully")
        
        return ' '.join(translated_chunks)
    
    except Exception as e:
        print(f"[Translation] Grok AI translation failed: {e}")
        return None


def extract_youtube_playlist_transcript(url: str) -> Tuple[str, str]:
    """
    Extract transcripts from all videos in a YouTube playlist.
    Supports multilingual videos with auto-translation.
    Uses pytube to get video IDs from playlist without API key.
    
    Returns:
        Tuple of (combined_transcript, primary_source_language)
    """
    try:
        import re
        from youtube_transcript_api import YouTubeTranscriptApi
        ytt_api = YouTubeTranscriptApi()
        
        # Extract playlist ID
        playlist_match = re.search(r'list=([^&]+)', url)
        if not playlist_match:
            raise ValueError("Invalid playlist URL")
        
        playlist_id = playlist_match.group(1)
        
        # Try to extract video IDs using pytube
        try:
            from pytube import Playlist
            
            playlist = Playlist(url)
            video_ids = []
            
            # Extract video IDs from playlist
            for video_url in playlist.video_urls:
                video_id = extract_video_id(video_url)
                if video_id:
                    video_ids.append(video_id)
            
            if not video_ids:
                raise ValueError("No videos found in playlist")
            
            # Extract transcripts from all videos (with multilingual support)
            all_transcripts = []
            successful_count = 0
            failed_count = 0
            detected_languages = []
            
            for video_id in video_ids:
                try:
                    text, lang = _extract_transcript_multilingual(video_id)
                    all_transcripts.append(text)
                    detected_languages.append(lang)
                    successful_count += 1
                except Exception as e:
                    print(f"Failed to get transcript for video {video_id}: {str(e)}")
                    failed_count += 1
                    continue
            
            if not all_transcripts:
                raise ValueError(f"Could not extract transcripts from any videos in playlist. {failed_count} videos failed.")
            
            # Determine primary language
            primary_lang = max(set(detected_languages), key=detected_languages.count) if detected_languages else 'en'
            
            # Combine all transcripts
            combined_transcript = '\n\n--- Next Video ---\n\n'.join(all_transcripts)
            
            # Add summary at the beginning
            lang_name = LANGUAGE_NAMES.get(primary_lang, primary_lang)
            summary = f"Playlist Transcript Summary: {successful_count} videos processed successfully"
            if failed_count > 0:
                summary += f", {failed_count} videos failed"
            if primary_lang != 'en':
                summary += f" (Original language: {lang_name}, auto-translated to English)"
            
            return f"{summary}\n\n{combined_transcript}", primary_lang
        
        except ImportError:
            # Fallback: Try using requests to scrape playlist page
            return extract_playlist_via_scraping(url, playlist_id)
    
    except Exception as e:
        error_msg = f"Playlist extraction failed: {str(e)}"
        print(error_msg)
        return f"Error: {error_msg}", 'en'


def extract_playlist_via_scraping(url: str, playlist_id: str) -> Tuple[str, str]:
    """
    Fallback method to extract playlist videos using web scraping.
    Supports multilingual transcripts.
    
    Returns:
        Tuple of (combined_transcript, primary_source_language)
    """
    try:
        import requests
        import re
        
        # Fetch playlist page
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        
        # Extract video IDs using regex
        video_id_pattern = r'"videoId":"([^"]+)"'
        video_ids = list(set(re.findall(video_id_pattern, response.text)))
        
        if not video_ids:
            raise ValueError("No videos found in playlist")
        
        # Limit to first 20 videos to avoid timeout
        video_ids = video_ids[:20]
        
        # Extract transcripts with multilingual support
        all_transcripts = []
        successful_count = 0
        failed_count = 0
        detected_languages = []
        
        for video_id in video_ids:
            try:
                text, lang = _extract_transcript_multilingual(video_id)
                all_transcripts.append(text)
                detected_languages.append(lang)
                successful_count += 1
            except Exception as e:
                failed_count += 1
                continue
        
        if not all_transcripts:
            raise ValueError(f"Could not extract transcripts from any videos. {failed_count} videos failed.")
        
        # Determine primary language
        primary_lang = max(set(detected_languages), key=detected_languages.count) if detected_languages else 'en'
        
        # Combine all transcripts
        combined_transcript = '\n\n--- Next Video ---\n\n'.join(all_transcripts)
        
        # Add summary
        lang_name = LANGUAGE_NAMES.get(primary_lang, primary_lang)
        summary = f"Playlist Transcript Summary: {successful_count} videos processed successfully"
        if failed_count > 0:
            summary += f", {failed_count} videos failed"
        if primary_lang != 'en':
            summary += f" (Original language: {lang_name}, auto-translated to English)"
        
        return f"{summary}\n\n{combined_transcript}", primary_lang
    
    except Exception as e:
        error_msg = f"Playlist scraping failed: {str(e)}"
        print(error_msg)
        return f"Error: {error_msg}. Please try installing pytube: pip install pytube", 'en'


def extract_video_id(url: str) -> str:
    """Extract video ID from YouTube URL"""
    patterns = [
        r'(?:youtube\.com\/watch\?v=|youtu\.be\/)([^&\n?#]+)',
        r'youtube\.com\/embed\/([^&\n?#]+)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    
    return None


def extract_pdf_text(file_path: str) -> str:
    """
    Extract text from PDF file with error handling
    """
    try:
        import PyPDF2
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")
        
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            if len(pdf_reader.pages) == 0:
                raise ValueError("PDF file is empty")
            
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        if not text.strip():
            raise ValueError("No text could be extracted from PDF")
        
        return text.strip()
    
    except FileNotFoundError as e:
        return f"Error: {str(e)}"
    except Exception as e:
        error_msg = f"PDF text extraction failed: {str(e)}"
        print(error_msg)
        return f"Error: {error_msg}. The PDF may be corrupted or image-based."


def extract_ppt_text(file_path: str) -> str:
    """
    Extract text from PowerPoint file with error handling
    """
    try:
        from pptx import Presentation
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            text += slide_text
        
        if not text.strip():
            raise ValueError("No text could be extracted from PowerPoint")
        
        return text.strip()
    
    except FileNotFoundError as e:
        return f"Error: {str(e)}"
    except Exception as e:
        error_msg = f"PPT text extraction failed: {str(e)}"
        print(error_msg)
        return f"Error: {error_msg}. The file may be corrupted."


def extract_word_text(file_path: str) -> str:
    """
    Extract text from Word document with error handling
    """
    try:
        import docx
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Word document not found: {file_path}")
        
        doc = docx.Document(file_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        if not text.strip():
            raise ValueError("No text could be extracted from Word document")
        
        return text.strip()
    
    except FileNotFoundError as e:
        return f"Error: {str(e)}"
    except Exception as e:
        error_msg = f"Word text extraction failed: {str(e)}"
        print(error_msg)
        return f"Error: {error_msg}. The file may be corrupted."


def extract_key_concepts(text: str, max_concepts: int = 10) -> List[str]:
    """
    Extract key concepts from text
    Simple keyword extraction based on frequency and importance
    """
    if not text:
        return []
    
    # Clean text
    text = text.lower()
    text = re.sub(r'[^\w\s]', ' ', text)
    
    # Remove common stop words
    stop_words = {
        'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
        'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were', 'be',
        'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will',
        'would', 'should', 'could', 'may', 'might', 'must', 'can', 'this',
        'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they',
        'what', 'which', 'who', 'when', 'where', 'why', 'how', 'all', 'each',
        'every', 'both', 'few', 'more', 'most', 'other', 'some', 'such', 'no',
        'nor', 'not', 'only', 'own', 'same', 'so', 'than', 'too', 'very'
    }
    
    # Split into words
    words = text.split()
    
    # Count word frequency (excluding stop words and short words)
    word_freq = {}
    for word in words:
        if len(word) > 3 and word not in stop_words:
            word_freq[word] = word_freq.get(word, 0) + 1
    
    # Sort by frequency
    sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
    
    # Get top concepts
    concepts = [word for word, freq in sorted_words[:max_concepts]]
    
    return concepts


def identify_concepts_from_questions(questions: List[Dict]) -> List[str]:
    """
    Extract concepts from generated questions
    """
    concepts = set()
    
    for q in questions:
        if 'concept' in q:
            concepts.add(q['concept'])
    
    return list(concepts)
